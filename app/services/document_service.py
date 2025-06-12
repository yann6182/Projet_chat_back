# app/services/document_service.py
from typing import Dict, List, BinaryIO, Optional
import os
import uuid
import docx2txt
import PyPDF2
import spacy
from pptx import Presentation
import language_tool_python
from spellchecker import SpellChecker
import re
import asyncio
import hashlib
from mistralai.client import MistralClient
from app.core.config import settings
from app.schemas.document import DocumentAnalysisResponse, SpellingError, GrammarError, LegalComplianceIssue
from app.services.chat_service import ChatService

class DocumentService:
    """
    Service pour l'analyse de documents.
    
    Ce service offre des fonctionnalités d'analyse orthographique, grammaticale et juridique
    des documents. La fonctionnalité RAG (Retrieval-Augmented Generation) est utilisée pour
    l'analyse de conformité juridique avancée.
    """
    def load_technical_terms(self):
        """
        Charge les termes techniques à ignorer lors de la vérification orthographique.
        Les termes peuvent être listés dans un fichier data/technical_terms.txt.
        """
        import os
        self.technical_terms = set()
        terms_file = os.path.join(os.path.dirname(__file__), '..', '..', 'data', 'technical_terms.txt')
        try:
            with open(terms_file, encoding='utf-8') as f:
                for line in f:
                    term = line.strip()
                    if term:
                        self.technical_terms.add(term.lower())
        except FileNotFoundError:
            pass
    
    def __init__(self, chat_service=None):
        # Services externes
        self.client = MistralClient(api_key=settings.MISTRAL_API_KEY)
        
        # Liste de termes techniques et mots fréquemment utilisés à ne pas signaler
        self.technical_terms = set()
        self.load_technical_terms()
        
        # Termes juridiques importants pour la vérification de base
        self.legal_terms = [
            "siret", "tva", "responsabilité", "clause", "contrat", 
            "obligation", "condition", "facturation", "paiement", "délai",
            "confidentialité", "propriété", "juridique", "légal", "statut"
        ]
        
        # Liste de mots vides à ignorer (articles, prépositions, etc.)
        self.stop_words = self._load_stop_words()
        
        # Initialiser les services d'IA
        self.model = "mistral-large-latest"  # Le modèle à utiliser
          # Initialiser le service de chat qui sera utilisé pour l'analyse RAG
        self.chat_service = chat_service
        self._setup_rag_service()
        
    def _setup_rag_service(self):
        """
        Vérifie et configure le service RAG s'il n'est pas déjà disponible
        """
        # Si le service de chat (RAG) n'est pas fourni, le créer
        if not self.chat_service:
            from app.services.chat_service import ChatService
            try:
                print("🔍 Tentative d'initialisation du service RAG pour l'analyse de documents...")
                self.chat_service = ChatService()
                
                # Vérifier si le service RAG est actif
                rag_available = hasattr(self.chat_service, 'use_chroma') and self.chat_service.use_chroma
                db_available = hasattr(self.chat_service, 'chroma_service') and self.chat_service.chroma_service
                
                if rag_available and db_available:
                    print("✅ Service RAG initialisé et prêt pour l'analyse juridique")
                else:
                    print("⚠️ Service RAG initialisé mais potentiellement incomplet:")
                    print(f"   - use_chroma: {rag_available}")
                    print(f"   - chroma_service: {db_available}")
            except Exception as e:
                print(f"❌ Erreur lors de l'initialisation du service RAG pour l'analyse de documents: {e}")
                # Même en cas d'erreur, on continue sans RAG
                self.chat_service = None
    
        # Charger le modèle spaCy pour l'analyse linguistique en français
        try:
            self.nlp = spacy.load("fr_core_news_md")
        except IOError:
            # Fallback vers le modèle plus petit si le grand n'est pas disponible
            self.nlp = spacy.load("fr_core_news_sm")
        
        self.upload_dir = "data/user_uploads"
        os.makedirs(self.upload_dir, exist_ok=True)
        
        # Initialiser les outils de vérification avancée
        self.grammar_tool = language_tool_python.LanguageTool('fr')
        self.spell_checker = SpellChecker(language='fr')
        
        # Cache pour les termes techniques à ne pas signaler comme erreurs
        self.technical_terms_cache = set([
            # Termes informatiques courants
            "web", "angular", "react", "vue", "node", "javascript", "typescript", "html",
            "css", "api", "php", "sql", "nosql", "mongodb", "firebase", "aws", "azure",
            "frontend", "backend", "fullstack", "bootstrap", "jquery", "github", "gitlab",
            # Termes liés aux Juniors-Entreprises
            "cnje", "junior-entreprise", "je", "tva", "siret", "ape", "crm", "erp",
            # Abréviations courantes
            "cf", "etc", "ex", "nb", "url", "www", "http", "https", "app", "sdk"
        ])
        
        # Liste de termes juridiques spécifiques aux Juniors Entreprises
        self.legal_terms = [
            "junior entreprise", "statut associatif", "CNJE", "étudiant entrepreneur",
            "prestation intellectuelle", "convention", "facturation", "TVA",
            "responsabilité civile professionnelle", "cotisation", "assemblée générale",
            "contrat de prestation", "devis", "facture", "TVA intracommunautaire",
            "responsabilité civile", "assurance", "statuts", "règlement intérieur",
            "conseil d'administration", "assemblée générale ordinaire", "AGO",
            "assemblée générale extraordinaire", "AGE", "bilan financier",
            "compte de résultat", "trésorier", "président", "vice-président",
            "secrétaire général", "commissaire aux comptes", "auditeur",
            "junior-entreprise", "JE", "étudiant", "école", "université",
            "formation", "compétences", "mission", "client", "prospect",
            "commercial", "développement", "qualité", "suivi", "livrable"
        ]
        
        # Expressions juridiques importantes à vérifier
        self.legal_expressions = [
            r"numéro de TVA",
            r"TVA intracommunautaire",
            r"responsabilité civile professionnelle",
            r"assurance responsabilité civile",
            r"numéro SIRET",
            r"code APE",
            r"conditions générales",
            r"clause de confidentialité",
            r"propriété intellectuelle",
            r"droit d'auteur",
            r"délai de paiement",
            r"pénalités de retard",
            r"tribunal compétent",
            r"droit applicable"
        ]
        
        # Dictionnaire de corrections courantes pour l'orthographe française
        self.common_corrections = {
            "contract": "contrat",
            "signature electronique": "signature électronique",
            "assemblee": "assemblée",
            "president": "président",
            "tresorier": "trésorier",
            "prestation": "prestation",
            "junior-enterprise": "junior-entreprise",
            "cnje": "CNJE",
            "developper": "développer",
            "etudiant": "étudiant",
            "universite": "université",
            "ecole": "école"
        }
    
        
    async def save_document(self, file: BinaryIO, filename: str) -> str:
        """
        Sauvegarde un document téléchargé et retourne son identifiant unique.
        
        Args:
            file: Le fichier téléchargé
            filename: Le nom du fichier
            
        Returns:
            L'identifiant unique du document
        """
        # Générer un ID unique pour le document
        document_id = str(uuid.uuid4())
        
        # Déterminer l'extension du fichier
        _, ext = os.path.splitext(filename)
        
        # Créer le chemin complet
        file_path = os.path.join(self.upload_dir, f"{document_id}{ext}")
        
        # Sauvegarder le fichier
        with open(file_path, "wb") as f:
            f.write(file.read())
            
        return document_id
        
    async def analyze_document(self, document_id: str) -> DocumentAnalysisResponse:
        """
        Analyse un document pour vérifier son orthographe, sa grammaire et sa conformité légale.
        
        Args:
            document_id: L'identifiant du document à analyser
            
        Returns:
            Le résultat de l'analyse du document
        """
        # Trouver le fichier correspondant à l'ID
        file_path, filename = self.find_document_by_id(document_id)
        
        if not file_path:
            raise FileNotFoundError(f"Document avec l'ID {document_id} non trouvé")
            
        # Extraire le texte du document
        text = self.extract_text(file_path)
        
        # Pour les documents très longs, optimiser le texte pour l'analyse
        if len(text) > 20000:  # Si plus de 20 000 caractères
            optimized_text = self._optimize_text_for_analysis(text)
            print(f"Document optimisé de {len(text)} à {len(optimized_text)} caractères pour l'analyse")
        else:
            optimized_text = text
          # Analyser le texte
        spelling_errors = await self.check_spelling(optimized_text)
        grammar_errors = self.check_grammar(optimized_text)
        legal_compliance_issues = await self.check_legal_compliance(optimized_text)
        
        # Calculer un score de conformité global
        total_issues = len(spelling_errors) + len(grammar_errors) + len(legal_compliance_issues)
        # Ajuster le calcul du score pour être plus équilibré
        max_allowed_issues = min(100, len(text) // 200)  # 1 erreur tous les 200 caractères maximum
        compliance_score = max(0.0, 1.0 - (total_issues / max_allowed_issues)) if total_issues > 0 else 1.0
        
        # Générer des suggestions d'amélioration
        suggestions = self.generate_suggestions(text, spelling_errors, grammar_errors, legal_compliance_issues)
        
        return DocumentAnalysisResponse(
            document_id=document_id,
            filename=filename,
            spelling_errors=spelling_errors,
            grammar_errors=grammar_errors,
            legal_compliance_issues=legal_compliance_issues,
            overall_compliance_score=compliance_score,
            suggestions=suggestions
        )
        
    def find_document_by_id(self, document_id: str) -> tuple:
        """
        Recherche un document par son identifiant.
        
        Args:
            document_id: L'identifiant du document
            
        Returns:
            Un tuple contenant le chemin du fichier et son nom
        """
        for filename in os.listdir(self.upload_dir):
            if filename.startswith(document_id):
                return os.path.join(self.upload_dir, filename), filename
        return None, None
        
    def extract_text(self, file_path: str) -> str:
        """
        Extrait le texte d'un document selon son format.
        
        Args:
            file_path: Le chemin du fichier
            
        Returns:
            Le texte extrait du document
        """
        _, ext = os.path.splitext(file_path)
        ext = ext.lower()
        
        if ext == ".pdf":
            return self.extract_text_from_pdf(file_path)
        elif ext in [".docx", ".doc"]:
            return self.extract_text_from_docx(file_path)
        elif ext in [".pptx", ".ppt"]:
            return self.extract_text_from_pptx(file_path)
        elif ext == ".txt":
            with open(file_path, "r", encoding="utf-8") as f:
                return f.read()
        else:
            raise ValueError(f"Format de fichier non pris en charge: {ext}")
    
    def extract_text_from_pdf(self, file_path: str) -> str:
        """
        Extrait le texte d'un fichier PDF.
        
        Args:
            file_path: Le chemin du fichier PDF
            
        Returns:
            Le texte extrait du PDF
        """
        text = ""
        with open(file_path, "rb") as f:
            pdf_reader = PyPDF2.PdfReader(f)
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
        return text
        
    def extract_text_from_docx(self, file_path: str) -> str:
        """
        Extrait le texte d'un fichier DOCX.
        
        Args:
            file_path: Le chemin du fichier DOCX
            
        Returns:
            Le texte extrait du DOCX
        """
        return docx2txt.process(file_path)
        
    def extract_text_from_pptx(self, file_path: str) -> str:
        """
        Extrait le texte d'un fichier PowerPoint.
        
        Args:
            file_path: Le chemin du fichier PowerPoint
            
        Returns:
            Le texte extrait du PowerPoint
        """
        text = ""
        prs = Presentation(file_path)
        def extract_text_from_pptx(self, file_path: str) -> str:
            text = ""
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    # on vérifie que la forme contient du texte
                    if getattr(shape, "has_text_frame", False):
                        text += shape.text + "\n"
            return text
    async def check_spelling(self, text: str) -> List[SpellingError]:
        """
        Vérifie l'orthographe du texte avec un correcteur orthographique avancé,
        en utilisant Mistral pour valider les termes techniques et étrangers.
        
        Args:
            text: Le texte à vérifier
            
        Returns:
            Une liste d'erreurs d'orthographe
        """
        errors = []
        
        # Si le texte est vide, retourner immédiatement
        if not text or len(text) < 10:
            return errors
            
        # Limiter la taille du texte pour les très longs documents
        max_chars = 20000
        if len(text) > max_chars:
            print(f"Texte trop long ({len(text)} caractères), analyse limitée aux {max_chars} premiers caractères.")
            text = text[:max_chars]
        
        # Tokeniser le texte en mots - optimisé pour inclure plus de caractères spéciaux utilisés en programmation
        words = re.findall(r'\b[a-zA-ZàâäéèêëïîôöùûüçÀÂÄÉÈÊËÏÎÔÖÙÛÜÇ\-_]+\b', text)
        
        # Pour optimiser les performances, on va regrouper les mots par lots
        batch_size = 200  # Traiter plus de mots par lots pour améliorer les performances
        unique_words = set(word for word in words if len(word) >= 3)  # Éliminer les doublons et mots courts
        
        # Créer une structure de données plus efficace pour retrouver les positions des mots
        word_positions = {}
        
        # Prétraitement : identifier les termes techniques en une seule passe
        technical_terms = set()
        for word in unique_words:
            if self._is_technical_term(word):
                technical_terms.add(word.lower())
        
        # Traiter les mots par lots pour réduire le nombre d'appels au correcteur orthographique
        word_batches = list(unique_words)
        
        # Filtrer les termes techniques et juridiques avant la vérification orthographique
        filtered_batch = [
            word for word in word_batches 
            if word.lower() not in technical_terms and
            not any(term.lower() == word.lower() for term in self.legal_terms)
        ]
            
        if not filtered_batch:
            return errors
                
        # Vérifier tous les mots du lot en une seule opération
        misspelled = self.spell_checker.unknown(filtered_batch)
        
        # Si aucune erreur n'est détectée, retourner directement
        if not misspelled:
            return errors
            
        # Utiliser Mistral pour valider les termes qui semblent être des erreurs
        # mais qui pourraient être des termes techniques, anglais, etc.
        mistral_validations = await self._validate_with_mistral(list(misspelled))
        
        # Traiter seulement les mots que Mistral considère comme invalides
        actual_misspelled = [word for word in misspelled if not mistral_validations.get(word, False)]
        
        # Pour chaque mot réellement mal orthographié
        for word in actual_misspelled:
            # Calculer position seulement pour les mots mal orthographiés (paresseux)
            if word not in word_positions:
                # Rechercher toutes les occurrences du mot dans le texte
                pattern = r'\b' + re.escape(word) + r'\b'
                for match in re.finditer(pattern, text):
                    position = {"start": match.start(), "end": match.end()}
                    
                    # Obtenir des suggestions de correction avec mise en cache pour des appels répétés
                    candidates = self.spell_checker.candidates(word.lower())
                    suggestions = list(candidates)[:3] if candidates else []
                    
                    # Vérifier avec le dictionnaire de corrections personnalisé
                    if word.lower() in self.common_corrections:
                        suggestions.insert(0, self.common_corrections[word.lower()])
                    
                    if suggestions:  # Seulement ajouter si on a des suggestions
                        errors.append(
                            SpellingError(
                                word=word,
                                position=position,
                                suggestions=suggestions
                            )
                        )
                    break  # Ne prendre que la première occurrence pour éviter les doublons
        
        return errors
        
    def check_grammar(self, text: str) -> List[GrammarError]:
        """
        Vérifie la grammaire du texte avec LanguageTool.
        
        Args:
            text: Le texte à vérifier
            
        Returns:
            Une liste d'erreurs grammaticales
        """
        errors = []
        
        # Si le texte est vide, retourner immédiatement
        if not text or len(text) < 10:
            return errors
        
        # Optimisation pour les textes longs - analyse par sections
        max_section_length = 5000  # Longueur maximale de chaque section
        
        # Estimer le temps d'analyse en fonction de la longueur du texte
        estimated_time_per_char = 0.0001  # Estimation très approximative
        estimated_time_seconds = len(text) * estimated_time_per_char
        
        # Approche adaptative en fonction de la longueur
        if len(text) <= max_section_length:
            # Pour les textes courts, analyser tout le document
            return self._check_grammar_section(text, 0)
        elif len(text) <= 15000:
            # Pour les textes moyens, analyser le début et la fin
            section_size = max_section_length // 2  # Moitié de la taille maximale
            
            # Analyser le début du document
            errors.extend(self._check_grammar_section(text[:section_size], 0))
            
            # Analyser la fin du document
            if len(text) > section_size + 500:  # Éviter le chevauchement
                start_pos = len(text) - section_size
                errors.extend(self._check_grammar_section(text[-section_size:], start_pos))
                
            print(f"Texte de taille moyenne ({len(text)} caractères), analyse limitée au début et à la fin ({2 * section_size} caractères).")
        else:
            # Pour les textes très longs, analyser début, milieu et des échantillons
            print(f"Texte long ({len(text)} caractères), analyse échantillonnée.")
            
            # 1. Début du document (40%)
            start_portion = int(max_section_length * 0.4)
            errors.extend(self._check_grammar_section(text[:start_portion], 0))
            
            # 2. Milieu du document (30%)
            mid_point = len(text) // 2
            mid_portion = int(max_section_length * 0.3) 
            mid_start = mid_point - (mid_portion // 2)
            mid_text = text[mid_start:mid_start + mid_portion]
            errors.extend(self._check_grammar_section(mid_text, mid_start))
            
            # 3. Fin du document (30%)
            end_portion = int(max_section_length * 0.3)
            end_text = text[-end_portion:]
            errors.extend(self._check_grammar_section(end_text, len(text) - end_portion))
            
            # 4. Échantillons aléatoires pour les très longs documents
            if len(text) > 50000:
                import random
                # Prendre quelques échantillons aléatoires de paragraphes
                paragraphs = text.split('\n\n')
                if len(paragraphs) > 5:
                    samples = random.sample(paragraphs, min(5, len(paragraphs) // 10))
                    for sample in samples:
                        if len(sample) > 100:  # Ignorer les paragraphes trop courts
                            # Trouver la position du paragraphe dans le texte original
                            sample_pos = text.find(sample)
                            if sample_pos != -1:
                                errors.extend(self._check_grammar_section(sample, sample_pos))
        
        return errors
    
    def _check_grammar_section(self, text_section: str, offset: int = 0) -> List[GrammarError]:
        """
        Vérifie la grammaire d'une section de texte.
        
        Args:
            text_section: La section de texte à vérifier
            offset: Position de début de la section dans le texte complet
            
        Returns:
            Liste des erreurs grammaticales
        """
        errors = []
        
        try:
            # Utiliser LanguageTool pour la vérification grammaticale
            # Configuration optimisée pour réduire les faux positifs
            matches = self.grammar_tool.check(
                text_section,
                # Désactiver certaines règles qui génèrent souvent des faux positifs
                disabled_rules=['WHITESPACE_RULE', 'UPPERCASE_SENTENCE_START']
            )
            
            # Construire un ensemble de termes techniques pour une recherche plus rapide
            technical_terms = set()
            words = re.findall(r'\b\w+\b', text_section)
            for word in set(words):
                if self._is_technical_term(word):
                    technical_terms.add(word.lower())
            
            for match in matches:
                # Filtrer les erreurs non pertinentes (orthographe déjà vérifiée ailleurs)
                if match.category in ['TYPOS', 'SPELLING', 'CASING']:
                    continue
                
                # Vérifier si le texte en erreur est un terme technique
                error_text = text_section[match.offset:match.offset + match.errorLength]
                
                # Ignorer les erreurs pour les termes techniques
                if error_text.lower() in technical_terms or self._is_technical_term(error_text):
                    continue
                    
                # Optimiser : Filtrer les cas où le match contient un terme technique
                contains_tech_term = any(term in error_text.lower() for term in technical_terms)
                if contains_tech_term:
                    continue
                
                # Ajuster la position pour correspondre au texte complet
                adjusted_start = match.offset + offset
                adjusted_end = adjusted_start + match.errorLength
                
                error = GrammarError(
                    text=error_text,
                    position={"start": adjusted_start, "end": adjusted_end},
                    message=match.message,
                    suggestions=match.replacements[:3]  # Limiter à 3 suggestions
                )
                errors.append(error)
                
        except Exception as e:
            print(f"Erreur lors de la vérification grammaticale de la section: {e}")
            # Fallback vers l'ancienne méthode si LanguageTool échoue
            try:
                fallback_errors = self._check_grammar_fallback(text_section)
                # Ajuster les positions des erreurs
                for error in fallback_errors:
                    error.position["start"] += offset
                    error.position["end"] += offset
                errors.extend(fallback_errors)
            except Exception:
                # En cas d'échec du fallback, simplement ignorer cette section
                pass
            
        return errors
    async def check_legal_compliance(self, text: str) -> List[LegalComplianceIssue]:
        """
        Vérifie la conformité légale du texte en analysant avec la base de connaissances.
        
        Args:
            text: Le texte à vérifier
            
        Returns:
            Une liste de problèmes de conformité légale
        """
        issues = []
        
        # 1. Vérifications automatiques de base
        issues.extend(self._check_basic_legal_requirements(text))
        
        # 2. Analyse avancée avec la base de connaissances (RAG) - RÉACTIVÉE avec optimisations
        try:
            # Pour les documents très longs, limiter l'analyse à 3000 caractères pertinents
            max_length = 3000
            if len(text) > max_length:
                # Prendre les 1500 premiers caractères (début du document)
                start_text = text[:1500]
                # Chercher des mots-clés juridiques importants
                keywords = ["responsabilité", "conditions générales", "TVA", "clause", "propriété", 
                           "contrat", "statut", "association", "CNJE", "junior-entreprise",
                           "siret", "facturation", "règlement", "assemblée", "cotisation"]
                legal_chunks = []
                
                # Parcourir les mots-clés et extraire le contexte autour d'eux
                for keyword in keywords:
                    if keyword in text.lower():
                        idx = text.lower().find(keyword)
                        if idx >= 0:
                            # Prendre 200 caractères autour du mot-clé
                            start_idx = max(0, idx - 100)
                            end_idx = min(len(text), idx + 100)
                            legal_chunks.append(text[start_idx:end_idx])
                
                # Combiner le début du document et les sections pertinentes
                text_to_analyze = start_text + "\n...\n" + "\n".join(legal_chunks)
                
                # Limiter à la taille maximale si nécessaire
                if len(text_to_analyze) > max_length:
                    text_to_analyze = text_to_analyze[:max_length]
                    
                legal_analysis = await self._analyze_with_knowledge_base(text_to_analyze)
            else:
                legal_analysis = await self._analyze_with_knowledge_base(text)
            
            issues.extend(legal_analysis)
            
        except Exception as e:
            print(f"Erreur lors de l'analyse juridique avancée: {e}")
            # En cas d'erreur, ajouter une notification
            issues.append(
                LegalComplianceIssue(
                    text="",
                    position={"start": 0, "end": 0},
                    issue_type="Erreur d'analyse RAG",
                    description="Une erreur s'est produite lors de l'analyse avancée avec la base de connaissances.",
                    recommendation="Veuillez vérifier l'état de la base de connaissances et réessayer."
                )
            )
        
        return issues
    
    def _check_basic_legal_requirements(self, text: str) -> List[LegalComplianceIssue]:
        """
        Vérifie les exigences légales de base.
        """
        issues = []
        text_lower = text.lower()
        
        # Vérifier la présence de termes juridiques importants
        missing_terms = []
        for term in self.legal_terms:
            if term not in text_lower:
                missing_terms.append(term)
        
        if len(missing_terms) > len(self.legal_terms) * 0.7:  # Si plus de 70% des termes manquent
            issues.append(
                LegalComplianceIssue(
                    text="",
                    position={"start": 0, "end": 0},
                    issue_type="Termes juridiques manquants",
                    description=f"Document incomplet : plusieurs termes juridiques importants sont absents. Exemples manquants : {', '.join(missing_terms[:5])}",
                    recommendation="Ajoutez les termes juridiques appropriés selon le type de document (contrat, statuts, etc.)."
                )
            )
        
        # Vérifications spécifiques selon le contenu
        if "tva" in text_lower and "numéro de tva" not in text_lower and "tva intracommunautaire" not in text_lower:
            tva_pos = text_lower.find("tva")
            issues.append(
                LegalComplianceIssue(
                    text="TVA",
                    position={"start": tva_pos, "end": tva_pos + 3},
                    issue_type="Mention légale incomplète",
                    description="La TVA est mentionnée mais le numéro de TVA intracommunautaire n'est pas précisé.",
                    recommendation="Ajoutez le numéro de TVA intracommunautaire de votre Junior-Entreprise."
                )
            )
        
        # Vérifier les mentions obligatoires pour les contrats
        if any(word in text_lower for word in ["contrat", "prestation", "service"]):
            contract_issues = self._check_contract_requirements(text)
            issues.extend(contract_issues)
        
        # Vérifier les statuts d'association
        if any(word in text_lower for word in ["statuts", "association", "assemblée"]):
            statute_issues = self._check_statute_requirements(text)
            issues.extend(statute_issues)
        
        return issues
    
    def _check_contract_requirements(self, text: str) -> List[LegalComplianceIssue]:
        """
        Vérifie les exigences spécifiques aux contrats.
        """
        issues = []
        text_lower = text.lower()
        
        required_clauses = [
            ("responsabilité civile", "clause de responsabilité civile"),
            ("conditions générales", "conditions générales de vente/prestation"),
            ("délai", "délais d'exécution"),
            ("paiement", "conditions de paiement"),
            ("propriété intellectuelle", "clause de propriété intellectuelle")
        ]
        
        for clause, description in required_clauses:
            if clause not in text_lower:
                issues.append(
                    LegalComplianceIssue(
                        text="",
                        position={"start": 0, "end": 0},
                        issue_type="Clause contractuelle manquante",
                        description=f"Clause manquante : {description}",
                        recommendation=f"Ajoutez une {description} dans votre contrat."
                    )
                )
        
        return issues
    
    def _check_statute_requirements(self, text: str) -> List[LegalComplianceIssue]:
        """
        Vérifie les exigences spécifiques aux statuts d'association.
        """
        issues = []
        text_lower = text.lower()
        
        required_elements = [
            ("objet social", "définition de l'objet social"),
            ("siège social", "adresse du siège social"),
            ("conseil d'administration", "composition du conseil d'administration"),
            ("assemblée générale", "modalités d'assemblée générale"),
            ("dissolution", "conditions de dissolution")
        ]
        
        for element, description in required_elements:
            if element not in text_lower:
                issues.append(
                    LegalComplianceIssue(
                        text="",
                        position={"start": 0, "end": 0},
                        issue_type="Élément statutaire manquant",
                        description=f"Élément manquant : {description}",
                        recommendation=f"Les statuts doivent inclure {description}."
                    )
                )
        
        return issues
    async def _analyze_with_knowledge_base(self, text: str) -> List[LegalComplianceIssue]:
        """
        Analyse le document avec la base de connaissances juridique.
        
        Cette méthode utilise RAG pour extraire des informations pertinentes
        de la base de connaissances juridique.
        """
        issues = []
        
        try:            # Préparer la requête pour l'analyse juridique avec demande EXPLICITE de citations
            legal_query = f"""
            Analysez ce document du point de vue juridique en vous basant UNIQUEMENT sur votre base de connaissances sur les Junior-Entreprises.
            Identifiez les éléments non conformes au droit français et aux spécificités des Junior-Entreprises.
            
            Document à analyser :
            {text[:2000]}
            
            Recherchez spécifiquement :
            1. Les mentions légales obligatoires manquantes
            2. Les clauses contractuelles non conformes
            3. Les éléments statutaires incorrects
            4. Les obligations fiscales et sociales non respectées
            
            INSTRUCTIONS STRICTES:
            - Pour CHAQUE problème identifié, vous DEVEZ suivre EXACTEMENT ce format en trois parties:
            
            PROBLÈME: [Description claire et précise du problème identifié]
            CITATION: [Citation EXACTE et TEXTUELLE provenant de votre base de connaissances juridique]
            RECOMMANDATION: [Votre recommandation concrète pour corriger le problème]
            
            - Ne JAMAIS omettre la section CITATION qui doit contenir du texte exact de votre base de connaissances
            - Si vous ne trouvez pas de citation pertinente dans votre base de connaissances, NE MENTIONNEZ PAS ce problème
            - Inclure au minimum 2 citations pertinentes de votre base de connaissances
            
            IMPORTANT: Je vais vérifier la présence des citations exactes. Si elles sont absentes, votre analyse sera considérée comme incomplète.
            """
              # Utiliser le service de chat pour analyser avec la base de connaissances
            from app.schemas.chat import ChatRequest
            import hashlib
            
            # Générer une clé de cache robuste basée sur le contenu du texte
            text_hash = hashlib.md5(text[:1500].encode('utf-8')).hexdigest()
            cache_key = f"legal_analysis_v2_{text_hash}"  # Ajouter un préfixe v2 pour forcer la mise à jour
            
            # Vérifier si la requête existe dans le cache pour éviter des requêtes répétitives
            from app.services.cache_service import get_cache
            cache = get_cache()
            cached_response = cache.get(cache_key)
            
            if cached_response and len(cached_response) > 200:  # Vérifier que la réponse est suffisamment longue
                # Si la réponse est en cache et semble valide, l'utiliser directement
                response_text = cached_response
                print(f"✅ Utilisation d'une réponse mise en cache pour l'analyse juridique ({len(response_text)} caractères)")
                
                # Vérifier si la réponse en cache contient des éléments au format structuré attendu
                if "PROBLÈME:" not in response_text.upper() and "CITATION:" not in response_text.upper():
                    print("⚠️ La réponse en cache ne semble pas contenir le format structuré attendu, force d'une nouvelle requête")
                    cached_response = None  # Forcer une nouvelle requête
            else:
                cached_response = None
                
            if not cached_response:
                print(f"🔄 Génération d'une nouvelle analyse juridique RAG")
                # Sinon, faire la requête
                chat_request = ChatRequest(query=legal_query)
                
                # Créer une conversation temporaire pour l'analyse
                import uuid
                temp_conversation_id = f"legal_analysis_{uuid.uuid4()}"
                
                try:
                    # Analyser avec le service de chat (qui utilise la base de connaissances)
                    response = await self.chat_service.process_query(
                        request=chat_request,
                        conversation_id=temp_conversation_id,
                        user_id=1  # Utilisateur système pour l'analyse
                    )
                    
                    if response and hasattr(response, 'response') and response.response:
                        response_text = response.response
                        print(f"✅ Réponse RAG obtenue ({len(response_text)} caractères)")
                        
                        # Vérifier la présence d'éléments au format attendu
                        if "PROBLÈME:" in response_text or "CITATION:" in response_text:
                            # Mettre en cache la réponse pour une utilisation future (TTL de 1 jour)
                            cache.set(cache_key, response_text, ttl=86400, persist=True)
                            print("💾 Réponse mise en cache pour utilisation future")
                        else:
                            print("⚠️ Format attendu non détecté dans la réponse")
                    else:
                        raise Exception("Pas de réponse du service de chat ou réponse invalide")
                except Exception as e:
                    # Capture de l'erreur spécifique pour un meilleur diagnostic
                    print(f"❌ Erreur lors de l'appel à process_query: {str(e)}")
                    # Essayer d'utiliser une réponse générique si disponible
                    response_text = "L'analyse juridique n'a pas pu être effectuée en raison d'un problème technique."
                    raise Exception(f"Échec de l'analyse RAG: {str(e)}")
            
            # Parser la réponse pour extraire les problèmes juridiques
            legal_issues = self._parse_legal_analysis_response(response_text)
            issues.extend(legal_issues)
        
        except Exception as e:
            print(f"Erreur lors de l'analyse avec la base de connaissances: {e}")
            # En cas d'erreur, ajouter un conseil général
            issues.append(
                LegalComplianceIssue(
                    text="",
                    position={"start": 0, "end": 0},
                    issue_type="Analyse juridique recommandée",
                    description="Une erreur s'est produite lors de l'analyse avec la base de connaissances.",
                    recommendation="Consultez un expert juridique ou votre référent CNJE pour valider la conformité de ce document."
                )
            )
        
        return issues
    def _parse_legal_analysis_response(self, response_text: str) -> List[LegalComplianceIssue]:
        """
        Parse la réponse de l'analyse juridique pour extraire les problèmes.
        Cette méthode améliorée recherche spécifiquement le format structuré avec citations.
        """
        issues = []
        
        # Format structuré avec extraction de bloc multilignes
        import re
        
        # Motif pour extraire des blocs complets "PROBLÈME:...CITATION:...RECOMMANDATION..."
        block_pattern = r"PROBLÈME\s*:\s*([\s\S]*?)(?:CITATION\s*:\s*([\s\S]*?))?(?:RECOMMANDATION\s*:\s*([\s\S]*?))?(?=PROBLÈME|$)"
        
        # Rechercher tous les blocs dans la réponse
        block_matches = re.finditer(block_pattern, response_text, re.IGNORECASE)
        
        for block_match in block_matches:
            # Extraire chaque partie du bloc
            problem_text = block_match.group(1).strip() if block_match.group(1) else ""
            citation_text = block_match.group(2).strip() if block_match.group(2) else ""
            recommendation_text = block_match.group(3).strip() if block_match.group(3) else ""
            
            # Si on n'a pas de citation dans le format structuré, essayer de l'extraire avec des patterns alternatifs
            if not citation_text:
                # Rechercher des citations entre guillemets dans le problème
                quote_pattern = r'[«"]([^«"]+)[»"]'
                quote_matches = re.findall(quote_pattern, problem_text)
                if quote_matches:
                    citation_text = "; ".join(quote_matches)
              # Si on a trouvé au moins un problème et une citation
            if problem_text:
                # Log de déboggage pour vérifier l'extraction
                print(f"Problème extrait: {problem_text[:50]}...")
                print(f"Citation extraite: {citation_text[:50]}..." if citation_text else "Pas de citation")
                
                # Ne créer l'issue que si on a une citation ou si on est en mode debug
                if citation_text:
                    issues.append(
                        LegalComplianceIssue(
                            text=citation_text,
                            position={"start": 0, "end": 0},
                            issue_type="Problème juridique identifié",
                            description=problem_text,
                            recommendation=recommendation_text or "Consultez un expert juridique pour plus d'informations."
                        )
                    )
        
        # Méthode améliorée pour extraire les citations si le format structuré a échoué
        if not issues:
            # Phase 1: Rechercher explicitement des paires problème-citation
            problem_citation_pattern = r'(?:problème|non-conformité|issue|erreur)[^\n]*?\:([^\n]+)(?:[^\n]*?citation[^\n]*?\:([^\n]+))'
            pc_matches = re.finditer(problem_citation_pattern, response_text, re.IGNORECASE)
            
            for pc_match in pc_matches:
                problem_desc = pc_match.group(1).strip()
                citation = pc_match.group(2).strip() if pc_match.group(2) else ""
                
                if citation:
                    issues.append(
                        LegalComplianceIssue(
                            text=citation,
                            position={"start": 0, "end": 0},
                            issue_type="Citation juridique",
                            description=problem_desc,
                            recommendation="Vérifiez la conformité avec les exigences juridiques."
                        )
                    )
            
            # Phase 2: Rechercher des citations entre guillemets si toujours pas de résultats
            if not issues:
                # Rechercher des mots-clés indiquant des problèmes juridiques
                problem_indicators = [
                    "manque", "manquant", "absent", "non conforme", "incorrecte", 
                    "obligatoire", "doit", "devrait", "nécessaire", "requis", "illégal", 
                    "non-respect", "violation", "interdit"
                ]
                
                # Rechercher des citations avec différents types de guillemets
                citation_pattern = r'[«"]([^"»]{15,})[»"]'
                citation_matches = re.finditer(citation_pattern, response_text)
                
                for match in citation_matches:
                    citation = match.group(1).strip()
                    # Trouver le contexte avant et après la citation
                    start_pos = max(0, match.start() - 150)
                    end_pos = min(len(response_text), match.end() + 150)
                    context = response_text[start_pos:end_pos]
                    
                    # Vérifier si le contexte contient un indicateur de problème
                    if any(indicator in context.lower() for indicator in problem_indicators):
                        # Extraire une description sensible du contexte
                        context_before = response_text[start_pos:match.start()].strip()
                        sentences_before = [s.strip() for s in context_before.split('.') if s.strip()]
                        description = sentences_before[-1] if sentences_before else "Problème juridique identifié"
                        
                        issues.append(
                            LegalComplianceIssue(
                                text=citation,
                                position={"start": 0, "end": 0},
                                issue_type="Citation juridique importante",
                                description=description,
                                recommendation="Vérifiez la conformité avec cette référence juridique."
                            )
                        )
              # Phase 3: Dernier recours - extraire des recommandations générales s'il n'y a toujours pas de citations
            if not issues:
                # Rechercher des phrases avec des recommandations
                recommendation_patterns = [
                    r"je (?:vous )?recommande[^.]+\.",
                    r"il (?:est|serait) (?:fortement |vivement )?recommandé[^.]+\.",
                    r"il (?:est|serait) (?:nécessaire|important|obligatoire)[^.]+\."
                ]
                
                for pattern in recommendation_patterns:
                    rec_matches = re.finditer(pattern, response_text, re.IGNORECASE)
                    for rec_match in rec_matches:
                        recommendation = rec_match.group(0).strip()
                        issues.append(
                            LegalComplianceIssue(
                                text="",
                                position={"start": 0, "end": 0},
                                issue_type="Recommandation juridique",
                                description="Analyse juridique: une recommandation a été identifiée",
                                recommendation=recommendation
                            )
                        )
                        
            # Si aucune citation n'a été trouvée malgré tous les efforts, ajouter un message d'erreur explicite
            if not issues:
                print("⚠️ Aucune citation juridique extraite de la réponse")
                issues.append(
                    LegalComplianceIssue(
                        text="",
                        position={"start": 0, "end": 0},
                        issue_type="Analyse juridique incomplète",
                        description="L'analyse n'a pas pu identifier de citations spécifiques de la base de connaissances juridique.",
                        recommendation="Essayez de reformuler ou d'enrichir le document pour une analyse plus précise."
                    )
                )
                
                # Log de la réponse complète pour le débogage
                print(f"Réponse RAG complète (tronquée): {response_text[:500]}...")
        
        return issues
        
    def generate_suggestions(self, text: str, spelling_errors: List[SpellingError], 
                            grammar_errors: List[GrammarError], 
                            legal_issues: List[LegalComplianceIssue]) -> List[str]:
        """
        Génère des suggestions d'amélioration basées sur les erreurs détectées.
        
        Args:
            text: Le texte analysé
            spelling_errors: Les erreurs d'orthographe détectées
            grammar_errors: Les erreurs grammaticales détectées
            legal_issues: Les problèmes de conformité légale détectés
            
        Returns:
            Une liste de suggestions d'amélioration
        """
        suggestions = []
        
        # Suggestions basées sur les erreurs d'orthographe
        for error in spelling_errors:
            if error.suggestions:
                suggestions.append(f"Orthographe : '{error.word}' pourrait être corrigé en {', '.join(error.suggestions)}")
        
        # Suggestions basées sur les erreurs grammaticales
        for error in grammar_errors:
            if error.suggestions:
                suggestions.append(f"Grammaire : {error.message} (ex: {', '.join(error.suggestions)})")
        
        # Suggestions basées sur les problèmes de conformité légale
        for issue in legal_issues:
            if issue.recommendation:
                suggestions.append(f"Conformité légale : {issue.recommendation}")
        
        # Suggestions générales d'amélioration
        if len(suggestions) == 0:
            suggestions.append("Aucune erreur détectée, document potentiellement conforme.")
        elif len(suggestions) <= 2:
            suggestions.append("Peu d'erreurs détectées, bon travail !")
        elif len(suggestions) <= 5:
            suggestions.append("Nombre modéré d'erreurs détectées, considérez les corrections suggérées.")
        else:
            suggestions.append("Nombre élevé d'erreurs détectées, une révision approfondie est recommandée.")
        
        return suggestions
    def _load_stop_words(self):
        """
        Charge les mots vides (stop words) en français à ignorer lors de l'analyse
        orthographique et lexicale.
        
        Returns:
            Un ensemble contenant les mots vides français
        """
        # Liste de base des mots vides français
        stop_words = {
            "le", "la", "les", "un", "une", "des", "du", "de", "ce", "cette", "ces",
            "mon", "ma", "mes", "ton", "ta", "tes", "son", "sa", "ses", "notre", "nos", "votre", "vos", "leur", "leurs",
            "et", "ou", "mais", "donc", "car", "ni", "or", "que", "quoi", "qui", "dont", "où",
            "à", "au", "aux", "avec", "chez", "dans", "de", "depuis", "derrière", "devant", "en", "entre", "jusque", "par", "pour", "sans", "sur", "vers",
            "je", "tu", "il", "elle", "on", "nous", "vous", "ils", "elles",
            "me", "te", "se", "lui", "y", "en",
            "être", "avoir", "faire", "dire", "aller", "voir", "pouvoir", "vouloir", "devoir", "falloir",
            "suis", "es", "est", "sommes", "êtes", "sont", "étais", "était", "étions", "étiez", "étaient",
            "ai", "as", "a", "avons", "avez", "ont", "avais", "avait", "avions", "aviez", "avaient",
            "fais", "fait", "faisons", "faites", "font", "faisais", "faisait", "faisions", "faisiez", "faisaient",
            "dis", "dit", "disons", "dites", "disent", "disais", "disait", "disions", "disiez", "disaient",
            "vais", "vas", "va", "allons", "allez", "vont", "allais", "allait", "allions", "alliez", "allaient",
            "vois", "voit", "voyons", "voyez", "voient", "voyais", "voyait", "voyions", "voyiez", "voyaient",
            "peux", "peut", "pouvons", "pouvez", "peuvent", "pouvais", "pouvait", "pouvions", "pouviez", "pouvaient",
            "veux", "veut", "voulons", "voulez", "veulent", "voulais", "voulait", "voulions", "vouliez", "voulaient",
            "dois", "doit", "devons", "devez", "doivent", "devais", "devait", "devions", "deviez", "devaient",
            "faut", "fallait"
        }
        
        # Essayer d'utiliser spaCy si disponible pour obtenir une liste plus complète
        try:
            if hasattr(self, 'nlp'):
                # Obtenir les stop words de spaCy
                spacy_stop_words = self.nlp.Defaults.stop_words
                stop_words.update(spacy_stop_words)
        except Exception:
            # En cas d'erreur, utiliser simplement la liste de base
            pass
            
        return stop_words    

    def _is_technical_term(self, word: str) -> bool:
        """
        Vérifie si un mot est un terme technique qui doit être ignoré lors de la vérification orthographique.
        
        Args:
            word: Le mot à vérifier
            
        Returns:
            True si le mot est un terme technique, False sinon
        """
        if not word:
            return False
            
        # Convertir en minuscule pour la comparaison
        word_lower = word.lower()
        
        # Vérifier si le mot est dans notre ensemble de termes techniques
        if word_lower in self.technical_terms:
            return True
        
        # Vérifier si le mot est dans notre cache de termes techniques
        if hasattr(self, 'technical_terms_cache') and word_lower in self.technical_terms_cache:
            return True
            
        # Vérifier les acronymes (mots en majuscules de 2 caractères ou plus)
        if word.isupper() and len(word) >= 2:
            return True
            
        # Vérifier les termes techniques communs qui incluent des chiffres ou caractères spéciaux
        if any(char.isdigit() for char in word) and any(char.isalpha() for char in word):
            return True
            
        # Vérifier si c'est un terme lié aux Junior-Entreprises
        je_terms = ["junior", "entreprise", "je", "cnje", "jeh", "urssaf", "siret", "ape", "tva"]
        if any(term in word_lower for term in je_terms):
            return True
            
        # Vérifier les mots composés avec tiret qui contiennent des termes techniques
        if '-' in word:
            parts = word_lower.split('-')
            return any(part in self.technical_terms for part in parts if part)
        return False
    
    async def _validate_with_mistral(self, words: List[str]) -> Dict[str, bool]:
        """
        Utilise l'API Mistral pour valider si des mots apparemment mal orthographiés
        sont en fait des termes techniques, des noms propres, des mots étrangers, etc.
        
        Args:
            words: Liste des mots à valider
            
        Returns:
            Dictionnaire avec les mots comme clés et des booléens comme valeurs
            (True si le mot est valide, False sinon)
        """
        if not words:
            return {}
            
        # Limiter le nombre de mots à valider pour éviter des requêtes trop longues
        max_words_per_batch = 20
        all_validations = {}
        
        # Traiter par lots pour les documents avec beaucoup d'erreurs potentielles
        for i in range(0, len(words), max_words_per_batch):
            batch = words[i:i+max_words_per_batch]
            
            try:
                if not hasattr(self, 'chat_service') or not self.chat_service:
                    # Si le service de chat n'est pas disponible, considérer tous les mots comme valides
                    # pour éviter des faux positifs
                    batch_validations = {word: True for word in batch}
                else:
                    # Créer le prompt pour Mistral
                    words_list = "\n".join([f"- {word}" for word in batch])
                    prompt = f"""
                    Voici une liste de mots qui ont été identifiés comme potentiellement mal orthographiés en français :
                    
                    {words_list}
                    
                    Certains peuvent être des termes techniques, des noms propres, des mots étrangers (notamment anglais), 
                    des acronymes ou des abréviations valides. Pour chacun, indique uniquement "VALIDE" si le mot est 
                    correct dans un contexte technique, juridique ou d'entreprise, ou "INVALIDE" s'il s'agit vraiment 
                    d'une faute d'orthographe en français. Format attendu:
                    
                    mot1: VALIDE/INVALIDE
                    mot2: VALIDE/INVALIDE
                    etc.
                    """
                    
                    # Appeler l'API Mistral via le service de chat
                    from app.schemas.chat import ChatRequest
                    response = await self.chat_service.process_query(
                        request=ChatRequest(query=prompt),
                        conversation_id=f"spelling_validation_{uuid.uuid4()}",
                        user_id=1  # Utilisateur système
                    )
                    
                    # Parser la réponse pour extraire les validations
                    if response and hasattr(response, 'answer'):
                        response_text = response.answer
                        batch_validations = {}
                        
                        # Parcourir chaque ligne pour trouver les validations
                        for line in response_text.split('\n'):
                            line = line.strip()
                            if not line or ':' not in line:
                                continue
                                
                            parts = line.split(':', 1)
                            if len(parts) != 2:
                                continue
                                
                            word = parts[0].strip()
                            validation = parts[1].strip().upper()
                            
                            # Vérifier si le mot est dans notre batch
                            matching_word = next((w for w in batch if w.lower() == word.lower()), None)
                            if matching_word:
                                batch_validations[matching_word] = 'VALIDE' in validation
                        
                        # Pour les mots non traités, les considérer comme valides par défaut
                        for word in batch:
                            if word not in batch_validations:
                                batch_validations[word] = True
                    else:
                        # En cas d'erreur, considérer tous les mots comme valides
                        batch_validations = {word: True for word in batch}
                
                all_validations.update(batch_validations)
                
            except Exception as e:
                print(f"Erreur lors de la validation avec Mistral: {str(e)}")
                # En cas d'erreur, considérer tous les mots du batch comme valides
                # pour éviter des faux positifs
                for word in batch:
                    all_validations[word] = True
        
        return all_validations